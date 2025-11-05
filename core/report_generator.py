"""
Report Generation Module for Training App
Generates PowerPoint, Word, and PDF reports with charts
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.units import inch
from datetime import date
import io
import pandas as pd
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import seaborn as sns
sns.set_style("whitegrid")


def create_chart_image(df, chart_type, title, **kwargs):
    """
    Creates a chart and returns it as a BytesIO object

    Args:
        df: pandas DataFrame with data
        chart_type: 'bar', 'pie', 'line', etc.
        title: Chart title
        **kwargs: Additional arguments for the chart

    Returns:
        BytesIO object containing PNG image
    """
    fig, ax = plt.subplots(figsize=(10, 6))
    show_labels = kwargs.get('show_labels', True)

    if chart_type == 'bar':
        x_col = kwargs.get('x', df.columns[0])
        y_col = kwargs.get('y', df.columns[1])
        bars = ax.bar(df[x_col], df[y_col], color='steelblue')
        ax.set_xlabel(x_col)
        ax.set_ylabel(y_col)
        plt.xticks(rotation=45, ha='right')

        # Add data labels
        if show_labels:
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{int(height)}',
                       ha='center', va='bottom', fontsize=9, fontweight='bold')

    elif chart_type == 'horizontal_bar':
        x_col = kwargs.get('x', df.columns[0])
        y_col = kwargs.get('y', df.columns[1])
        bars = ax.barh(df[x_col], df[y_col], color='steelblue')
        ax.set_xlabel(y_col)
        ax.set_ylabel(x_col)

        # Add data labels
        if show_labels:
            for bar in bars:
                width = bar.get_width()
                ax.text(width, bar.get_y() + bar.get_height()/2.,
                       f'{width:.1f}%' if width < 100 else f'{int(width)}',
                       ha='left', va='center', fontsize=9, fontweight='bold',
                       bbox=dict(boxstyle='round,pad=0.3', facecolor='white', alpha=0.7))

    elif chart_type == 'pie':
        labels = df[df.columns[0]]
        values = df[df.columns[1]]
        colors_pie = ['#00CC96', '#EF553B', '#636EFA', '#FFA15A', '#19D3F3']
        ax.pie(values, labels=labels, autopct='%1.1f%%', colors=colors_pie, startangle=90)
        ax.axis('equal')

    elif chart_type == 'line':
        x_col = kwargs.get('x', df.columns[0])
        y_cols = kwargs.get('y', [df.columns[1]])
        for y_col in y_cols:
            line = ax.plot(df[x_col], df[y_col], marker='o', label=y_col, linewidth=2)

            # Add data labels
            if show_labels:
                for i, (x, y) in enumerate(zip(df[x_col], df[y_col])):
                    ax.text(i, y, f'{int(y)}', ha='center', va='bottom',
                           fontsize=8, fontweight='bold')

        ax.set_xlabel(x_col)
        ax.legend()
        plt.xticks(rotation=45, ha='right')

    elif chart_type == 'grouped_bar':
        x_col = kwargs.get('x', df.columns[0])
        y_cols = kwargs.get('y', df.columns[1:])
        bars_container = df.plot(kind='bar', x=x_col, y=list(y_cols), ax=ax)
        plt.xticks(rotation=45, ha='right')
        ax.legend()

        # Add data labels for grouped bars
        if show_labels:
            for container in ax.containers:
                ax.bar_label(container, fmt='%d', fontsize=8, fontweight='bold')

    ax.set_title(title, fontsize=14, fontweight='bold')
    ax.grid(True, alpha=0.3)
    plt.tight_layout()

    # Save to BytesIO
    img_io = io.BytesIO()
    plt.savefig(img_io, format='PNG', dpi=150, bbox_inches='tight')
    img_io.seek(0)
    plt.close(fig)

    return img_io


def create_pptx_report(report_data: dict) -> bytes:
    """
    Creates a PowerPoint presentation from comprehensive report data.

    Args:
        report_data: Dictionary containing all report sections

    Returns:
        bytes: PowerPoint file as bytes
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "Food Concepts Training Report"
    subtitle.text = f"Comprehensive Training Analytics\nGenerated: {date.today().strftime('%B %d, %Y')}"

    # Apply styling to title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Executive Summary Slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "Executive Summary"

    body = summary_slide.placeholders[1].text_frame
    body.clear()

    summary_data = report_data.get('summary', {})
    p = body.add_paragraph()
    p.text = f"Report Year: {report_data.get('year', date.today().year)}"
    p.font.size = Pt(18)
    p.font.bold = True

    body.add_paragraph()
    metrics = [
        f"📊 Total Trainings: {summary_data.get('total_trainings', 0)}",
        f"👥 Total Staff (All): {summary_data.get('total_staff', 0)}",
        f"✅ Active Staff: {summary_data.get('active_staff', 0)}",
        f"📚 Staff Trained (Active): {summary_data.get('staff_trained', 0)}",
        f"🏪 Stores Trained: {summary_data.get('stores_trained', 0)}",
        f"⏰ Pending Trainings: {summary_data.get('pending_trainings', 0)}",
    ]

    for metric in metrics:
        p = body.add_paragraph()
        p.text = metric
        p.level = 0
        p.font.size = Pt(16)

    # Training Frequency Slide
    if 'training_frequency' in report_data and not report_data['training_frequency'].empty:
        freq_slide = prs.slides.add_slide(prs.slide_layouts[5])
        freq_slide.shapes.title.text = "Training Frequency"

        # Add table
        freq_df = report_data['training_frequency']
        rows, cols = len(freq_df) + 1, 3  # Header + data rows, 3 columns

        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(4.5)

        table = freq_slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column headers
        table.cell(0, 0).text = "Training Title"
        table.cell(0, 1).text = "Times Held"
        table.cell(0, 2).text = "Date Range"

        # Style header row
        for col in range(cols):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)

        # Fill data
        for idx, row_data in enumerate(freq_df.itertuples(), start=1):
            table.cell(idx, 0).text = str(row_data[1])  # Training Title
            table.cell(idx, 1).text = str(row_data[2])  # Times Held
            date_range = f"{row_data[4]} to {row_data[5]}" if len(row_data) > 5 else ""
            table.cell(idx, 2).text = date_range

            for col in range(cols):
                table.cell(idx, col).text_frame.paragraphs[0].font.size = Pt(10)

    # Stores Coverage Slide
    if 'stores_coverage' in report_data:
        stores_slide = prs.slides.add_slide(prs.slide_layouts[1])
        stores_slide.shapes.title.text = "Stores Training Coverage"

        body = stores_slide.placeholders[1].text_frame
        body.clear()

        stores_data = report_data['stores_coverage']

        p = body.add_paragraph()
        p.text = f"Total Stores: {stores_data.get('total', 0)}"
        p.font.size = Pt(20)
        p.font.bold = True

        body.add_paragraph()

        p = body.add_paragraph()
        p.text = f"✅ Stores Trained: {stores_data.get('trained', 0)} ({stores_data.get('trained_pct', 0):.1f}%)"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 128, 0)

        p = body.add_paragraph()
        p.text = f"⏳ Stores Pending: {stores_data.get('pending', 0)} ({stores_data.get('pending_pct', 0):.1f}%)"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 69, 0)

    # Staff Participation Slide
    if 'staff_participation' in report_data:
        staff_slide = prs.slides.add_slide(prs.slide_layouts[1])
        staff_slide.shapes.title.text = "Staff Participation Analysis"

        body = staff_slide.placeholders[1].text_frame
        body.clear()

        staff_data = report_data['staff_participation']

        metrics = [
            (f"Total Eligible Staff: {staff_data.get('eligible', 0)}", 20, False),
            (f"Participated: {staff_data.get('participated', 0)} ({staff_data.get('participation_rate', 0):.1f}%)", 18, True),
            (f"Not Participated: {staff_data.get('not_participated', 0)}", 18, False),
        ]

        for text, size, bold in metrics:
            p = body.add_paragraph()
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            body.add_paragraph()

    # Food Safety Report Slide
    if 'food_safety' in report_data:
        food_slide = prs.slides.add_slide(prs.slide_layouts[1])
        food_slide.shapes.title.text = "Food Safety Compliance"

        body = food_slide.placeholders[1].text_frame
        body.clear()

        food_data = report_data['food_safety']

        p = body.add_paragraph()
        p.text = "Food Handlers / Food Safety Training Status"
        p.font.size = Pt(18)
        p.font.bold = True

        body.add_paragraph()

        metrics = [
            f"Required to be Tested: {food_data.get('required', 0)}",
            f"Already Tested: {food_data.get('tested', 0)} ({food_data.get('tested_pct', 0):.1f}%)",
            f"Never Tested: {food_data.get('never_tested', 0)}",
            f"Due for Testing: {food_data.get('due_count', 0)}",
        ]

        for metric in metrics:
            p = body.add_paragraph()
            p.text = metric
            p.font.size = Pt(16)

    # Compliance Dashboard Slide with Chart
    if 'compliance' in report_data and report_data['compliance']:
        compliance_slide = prs.slides.add_slide(prs.slide_layouts[5])
        compliance_slide.shapes.title.text = "Compliance Dashboard - All Required Trainings"

        compliance_df = pd.DataFrame(report_data['compliance'])

        # Add chart
        chart_img = create_chart_image(
            compliance_df[['Training', 'Compliance Rate']],
            'horizontal_bar',
            'Compliance Rate by Training Type',
            x='Training',
            y='Compliance Rate'
        )

        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5)
        compliance_slide.shapes.add_picture(chart_img, left, top, height=height)

    # Venue Analysis Slide with Chart
    if 'venue_analysis' in report_data and not report_data['venue_analysis'].empty:
        venue_slide = prs.slides.add_slide(prs.slide_layouts[5])
        venue_slide.shapes.title.text = "Training by Venue"

        venue_df = report_data['venue_analysis']

        chart_img = create_chart_image(
            venue_df[['Venue', 'Sessions']].head(10),
            'bar',
            'Top 10 Training Venues',
            x='Venue',
            y='Sessions'
        )

        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5)
        venue_slide.shapes.add_picture(chart_img, left, top, height=height)

    # Department Analysis Slide with Chart
    if 'department_coverage' in report_data and not report_data['department_coverage'].empty:
        dept_df = report_data['department_coverage']
        # Filter out rows with zero or negative values
        dept_df = dept_df[(dept_df['Total Staff'] > 0) & (dept_df['Total Trainings'] > 0)]

        if not dept_df.empty:
            dept_slide = prs.slides.add_slide(prs.slide_layouts[5])
            dept_slide.shapes.title.text = "Training Coverage by Department"

            chart_img = create_chart_image(
                dept_df[['Department', 'Total Staff', 'Trained Staff']],
                'grouped_bar',
                'Department Coverage Analysis',
                x='Department',
                y=['Total Staff', 'Trained Staff']
            )

            left = Inches(0.5)
            top = Inches(1.5)
            height = Inches(5)
            dept_slide.shapes.add_picture(chart_img, left, top, height=height)

    # Training by Department Slide with Bar Chart
    if 'department_coverage' in report_data and not report_data['department_coverage'].empty:
        dept_df = report_data['department_coverage']
        # Filter out rows with zero or negative values
        dept_df = dept_df[(dept_df['Total Staff'] > 0) & (dept_df['Total Trainings'] > 0)]

        if not dept_df.empty and 'Total Trainings' in dept_df.columns:
            dept_trainings_slide = prs.slides.add_slide(prs.slide_layouts[5])
            dept_trainings_slide.shapes.title.text = "Trainings by Department"

            chart_img = create_chart_image(
                dept_df[['Department', 'Total Trainings']].sort_values('Total Trainings', ascending=False),
                'bar',
                'Total Trainings by Department',
                x='Department',
                y='Total Trainings'
            )

            left = Inches(1)
            top = Inches(1.5)
            height = Inches(5)
            dept_trainings_slide.shapes.add_picture(chart_img, left, top, height=height)

    # Regional Coverage Slide with Chart
    if 'regional_coverage' in report_data and not report_data['regional_coverage'].empty:
        region_df = report_data['regional_coverage']
        # Filter out rows with zero or negative values
        region_df = region_df[(region_df['Total Staff'] > 0) & (region_df['Total Trainings'] > 0)]

        if not region_df.empty:
            regional_slide = prs.slides.add_slide(prs.slide_layouts[5])
            regional_slide.shapes.title.text = "Regional Training Coverage"

            chart_img = create_chart_image(
                region_df[['Region', 'Total Staff', 'Trained Staff']],
                'grouped_bar',
                'Regional Coverage Analysis',
                x='Region',
                y=['Total Staff', 'Trained Staff']
            )

            left = Inches(0.5)
            top = Inches(1.5)
            height = Inches(5)
            regional_slide.shapes.add_picture(chart_img, left, top, height=height)

    # Training by Region Slide with Bar Chart
    if 'regional_coverage' in report_data and not report_data['regional_coverage'].empty:
        region_df = report_data['regional_coverage']
        # Filter out rows with zero or negative values
        region_df = region_df[(region_df['Total Staff'] > 0) & (region_df['Total Trainings'] > 0)]

        if not region_df.empty and 'Total Trainings' in region_df.columns:
            region_trainings_slide = prs.slides.add_slide(prs.slide_layouts[5])
            region_trainings_slide.shapes.title.text = "Trainings by Region"

            chart_img = create_chart_image(
                region_df[['Region', 'Total Trainings']].sort_values('Total Trainings', ascending=False),
                'bar',
                'Total Trainings by Region',
                x='Region',
                y='Total Trainings'
            )

            left = Inches(1)
            top = Inches(1.5)
            height = Inches(5)
            region_trainings_slide.shapes.add_picture(chart_img, left, top, height=height)

    # Monthly Trends Slide with Chart
    if 'monthly_trends' in report_data and not report_data['monthly_trends'].empty:
        monthly_slide = prs.slides.add_slide(prs.slide_layouts[5])
        monthly_slide.shapes.title.text = f"Monthly Training Trends ({report_data.get('year', date.today().year)})"

        monthly_df = report_data['monthly_trends']

        chart_img = create_chart_image(
            monthly_df[['Month Name', 'Trainings', 'Unique Staff']],
            'line',
            'Monthly Training Activity',
            x='Month Name',
            y=['Trainings', 'Unique Staff']
        )

        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5)
        monthly_slide.shapes.add_picture(chart_img, left, top, height=height)

    # Pending Staff Slide
    if 'staff_yet_to_train' in report_data and not report_data['staff_yet_to_train'].empty:
        pending_slide = prs.slides.add_slide(prs.slide_layouts[1])
        pending_slide.shapes.title.text = "Staff Yet to be Trained"

        body = pending_slide.placeholders[1].text_frame
        body.clear()

        yet_df = report_data['staff_yet_to_train']

        p = body.add_paragraph()
        p.text = f"Total Staff Needing Training: {len(yet_df)}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 69, 0)

        body.add_paragraph()

        p = body.add_paragraph()
        p.text = "These staff members are missing required trainings and are not on probation."
        p.font.size = Pt(14)

        # Top departments needing training
        if 'department' in yet_df.columns:
            dept_counts = yet_df['department'].value_counts().head(5)
            body.add_paragraph()
            p = body.add_paragraph()
            p.text = "Top Departments:"
            p.font.size = Pt(16)
            p.font.bold = True

            for dept, count in dept_counts.items():
                p = body.add_paragraph()
                p.text = f"  • {dept}: {count} staff"
                p.font.size = Pt(14)
                p.level = 1

    # Save to bytes
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()


def create_word_report(report_data: dict) -> bytes:
    """
    Creates a Word document from comprehensive report data.

    Args:
        report_data: Dictionary containing all report sections

    Returns:
        bytes: Word document as bytes
    """
    doc = Document()

    # Title
    title = doc.add_heading('Food Concepts Training Report', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = doc.add_paragraph(f'Comprehensive Training Analytics\nGenerated: {date.today().strftime("%B %d, %Y")}')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # Executive Summary
    doc.add_heading('Executive Summary', 1)

    summary_data = report_data.get('summary', {})
    doc.add_paragraph(f"Report Year: {report_data.get('year', date.today().year)}")

    summary_table = doc.add_table(rows=7, cols=2)
    summary_table.style = 'Light Grid Accent 1'

    metrics = [
        ('Total Trainings', summary_data.get('total_trainings', 0)),
        ('Total Staff (All)', summary_data.get('total_staff', 0)),
        ('Active Staff', summary_data.get('active_staff', 0)),
        ('Staff Trained (Active)', summary_data.get('staff_trained', 0)),
        ('Stores Trained', summary_data.get('stores_trained', 0)),
        ('Pending Trainings', summary_data.get('pending_trainings', 0)),
    ]

    for idx, (label, value) in enumerate(metrics):
        summary_table.rows[idx].cells[0].text = label
        summary_table.rows[idx].cells[1].text = str(value)

    doc.add_page_break()

    # Training Frequency
    if 'training_frequency' in report_data and not report_data['training_frequency'].empty:
        doc.add_heading('Training Frequency', 1)

        freq_df = report_data['training_frequency']

        freq_table = doc.add_table(rows=len(freq_df) + 1, cols=4)
        freq_table.style = 'Light Grid Accent 1'

        # Headers
        headers = ['Training Title', 'Times Held', 'First Date', 'Last Date']
        for idx, header in enumerate(headers):
            freq_table.rows[0].cells[idx].text = header

        # Data
        for row_idx, row_data in enumerate(freq_df.itertuples(), start=1):
            freq_table.rows[row_idx].cells[0].text = str(row_data[1])
            freq_table.rows[row_idx].cells[1].text = str(row_data[2])
            freq_table.rows[row_idx].cells[2].text = str(row_data[4]) if len(row_data) > 4 else ""
            freq_table.rows[row_idx].cells[3].text = str(row_data[5]) if len(row_data) > 5 else ""

        doc.add_page_break()

    # Stores Coverage
    if 'stores_coverage' in report_data:
        doc.add_heading('Stores Training Coverage', 1)

        stores_data = report_data['stores_coverage']

        stores_table = doc.add_table(rows=3, cols=2)
        stores_table.style = 'Light Grid Accent 1'

        stores_table.rows[0].cells[0].text = 'Total Stores'
        stores_table.rows[0].cells[1].text = str(stores_data.get('total', 0))

        stores_table.rows[1].cells[0].text = 'Stores Trained'
        stores_table.rows[1].cells[1].text = f"{stores_data.get('trained', 0)} ({stores_data.get('trained_pct', 0):.1f}%)"

        stores_table.rows[2].cells[0].text = 'Stores Pending'
        stores_table.rows[2].cells[1].text = f"{stores_data.get('pending', 0)} ({stores_data.get('pending_pct', 0):.1f}%)"

        doc.add_page_break()

    # Staff Participation
    if 'staff_participation' in report_data:
        doc.add_heading('Staff Participation Analysis', 1)

        staff_data = report_data['staff_participation']

        staff_table = doc.add_table(rows=4, cols=2)
        staff_table.style = 'Light Grid Accent 1'

        staff_table.rows[0].cells[0].text = 'Eligible Staff'
        staff_table.rows[0].cells[1].text = str(staff_data.get('eligible', 0))

        staff_table.rows[1].cells[0].text = 'Participated'
        staff_table.rows[1].cells[1].text = f"{staff_data.get('participated', 0)} ({staff_data.get('participation_rate', 0):.1f}%)"

        staff_table.rows[2].cells[0].text = 'Not Participated'
        staff_table.rows[2].cells[1].text = str(staff_data.get('not_participated', 0))

        staff_table.rows[3].cells[0].text = 'Participation Rate'
        staff_table.rows[3].cells[1].text = f"{staff_data.get('participation_rate', 0):.1f}%"

        doc.add_page_break()

    # Food Safety Report
    if 'food_safety' in report_data:
        doc.add_heading('Food Safety Compliance Report', 1)

        food_data = report_data['food_safety']

        food_table = doc.add_table(rows=4, cols=2)
        food_table.style = 'Light Grid Accent 1'

        food_table.rows[0].cells[0].text = 'Required to be Tested'
        food_table.rows[0].cells[1].text = str(food_data.get('required', 0))

        food_table.rows[1].cells[0].text = 'Already Tested'
        food_table.rows[1].cells[1].text = f"{food_data.get('tested', 0)} ({food_data.get('tested_pct', 0):.1f}%)"

        food_table.rows[2].cells[0].text = 'Never Tested'
        food_table.rows[2].cells[1].text = str(food_data.get('never_tested', 0))

        food_table.rows[3].cells[0].text = 'Due for Testing'
        food_table.rows[3].cells[1].text = str(food_data.get('due_count', 0))

        if food_data.get('due_details'):
            doc.add_paragraph()
            doc.add_heading('Staff Due for Testing', 2)

            due_df = pd.DataFrame(food_data['due_details'])

            if len(due_df) > 0:
                due_table = doc.add_table(rows=len(due_df) + 1, cols=4)
                due_table.style = 'Light Grid Accent 1'

                # Headers
                due_table.rows[0].cells[0].text = 'Name'
                due_table.rows[0].cells[1].text = 'Department'
                due_table.rows[0].cells[2].text = 'Due Date'
                due_table.rows[0].cells[3].text = 'Status'

                # Data
                for row_idx, row_data in enumerate(due_df.itertuples(), start=1):
                    due_table.rows[row_idx].cells[0].text = str(row_data.name)
                    due_table.rows[row_idx].cells[1].text = str(row_data.department)
                    due_table.rows[row_idx].cells[2].text = str(row_data.due_date)
                    due_table.rows[row_idx].cells[3].text = str(row_data.status)

    doc.add_page_break()

    # Compliance Dashboard
    if 'compliance' in report_data and report_data['compliance']:
        doc.add_heading('Compliance Dashboard - All Required Trainings', 1)

        compliance_df = pd.DataFrame(report_data['compliance'])

        compliance_table = doc.add_table(rows=len(compliance_df) + 1, cols=4)
        compliance_table.style = 'Light Grid Accent 1'

        # Headers
        compliance_table.rows[0].cells[0].text = 'Training'
        compliance_table.rows[0].cells[1].text = 'Trained Staff'
        compliance_table.rows[0].cells[2].text = 'Compliance Rate (%)'
        compliance_table.rows[0].cells[3].text = 'Overdue'

        # Data
        for row_idx, row_data in enumerate(compliance_df.itertuples(), start=1):
            compliance_table.rows[row_idx].cells[0].text = str(row_data.Training)
            compliance_table.rows[row_idx].cells[1].text = str(row_data._2)  # Trained Staff
            compliance_table.rows[row_idx].cells[2].text = f"{row_data._3:.1f}"  # Compliance Rate
            compliance_table.rows[row_idx].cells[3].text = str(row_data.Overdue)

        # Add chart
        chart_img = create_chart_image(
            compliance_df[['Training', 'Compliance Rate']],
            'horizontal_bar',
            'Compliance Rate by Training Type',
            x='Training', y='Compliance Rate'
        )
        doc.add_paragraph()
        doc.add_picture(chart_img, width=DocxInches(6))

    doc.add_page_break()

    # Venue Analysis
    if 'venue_analysis' in report_data and not report_data['venue_analysis'].empty:
        doc.add_heading('Training by Venue Analysis', 1)

        venue_df = report_data['venue_analysis']

        venue_table = doc.add_table(rows=len(venue_df) + 1, cols=3)
        venue_table.style = 'Light Grid Accent 1'

        venue_table.rows[0].cells[0].text = 'Venue'
        venue_table.rows[0].cells[1].text = 'Sessions'
        venue_table.rows[0].cells[2].text = 'Unique Staff'

        for row_idx, row_data in enumerate(venue_df.itertuples(), start=1):
            venue_table.rows[row_idx].cells[0].text = str(row_data.Venue)
            venue_table.rows[row_idx].cells[1].text = str(row_data.Sessions)
            venue_table.rows[row_idx].cells[2].text = str(row_data._3)  # Unique Staff

        # Add chart
        chart_img = create_chart_image(
            venue_df[['Venue', 'Sessions']].head(10),
            'bar',
            'Top 10 Training Venues',
            x='Venue', y='Sessions'
        )
        doc.add_paragraph()
        doc.add_picture(chart_img, width=DocxInches(6))

    doc.add_page_break()

    # Department Coverage
    if 'department_coverage' in report_data and not report_data['department_coverage'].empty:
        dept_df = report_data['department_coverage']
        # Filter out rows with zero or negative values
        dept_df = dept_df[(dept_df['Total Staff'] > 0) & (dept_df['Total Trainings'] > 0)]

        if not dept_df.empty:
            doc.add_heading('Department Training Coverage', 1)

            dept_table = doc.add_table(rows=len(dept_df) + 1, cols=4)
            dept_table.style = 'Light Grid Accent 1'

            dept_table.rows[0].cells[0].text = 'Department'
            dept_table.rows[0].cells[1].text = 'Total Staff'
            dept_table.rows[0].cells[2].text = 'Trained Staff'
            dept_table.rows[0].cells[3].text = 'Coverage Rate (%)'

            for row_idx, row_data in enumerate(dept_df.itertuples(), start=1):
                dept_table.rows[row_idx].cells[0].text = str(row_data.Department)
                dept_table.rows[row_idx].cells[1].text = str(row_data._2)  # Total Staff
                dept_table.rows[row_idx].cells[2].text = str(row_data._3)  # Trained Staff
                if hasattr(row_data, '_5'):
                    dept_table.rows[row_idx].cells[3].text = f"{row_data._5:.1f}"  # Coverage Rate
                else:
                    dept_table.rows[row_idx].cells[3].text = "N/A"

            # Add chart
            chart_img = create_chart_image(
                dept_df[['Department', 'Total Staff', 'Trained Staff']],
                'grouped_bar',
                'Department Coverage Analysis',
                x='Department', y=['Total Staff', 'Trained Staff']
            )
            doc.add_paragraph()
            doc.add_picture(chart_img, width=DocxInches(6))

            # Add Training by Department bar chart
            if 'Total Trainings' in dept_df.columns:
                doc.add_paragraph()
                doc.add_heading('Trainings by Department', 2)
                chart_img2 = create_chart_image(
                    dept_df[['Department', 'Total Trainings']].sort_values('Total Trainings', ascending=False),
                    'bar',
                    'Total Trainings by Department',
                    x='Department', y='Total Trainings'
                )
                doc.add_picture(chart_img2, width=DocxInches(6))

    doc.add_page_break()

    # Regional Coverage
    if 'regional_coverage' in report_data and not report_data['regional_coverage'].empty:
        region_df = report_data['regional_coverage']
        # Filter out rows with zero or negative values
        region_df = region_df[(region_df['Total Staff'] > 0) & (region_df['Total Trainings'] > 0)]

        if not region_df.empty:
            doc.add_heading('Regional Training Coverage', 1)

            region_table = doc.add_table(rows=len(region_df) + 1, cols=3)
            region_table.style = 'Light Grid Accent 1'

            region_table.rows[0].cells[0].text = 'Region'
            region_table.rows[0].cells[1].text = 'Total Staff'
            region_table.rows[0].cells[2].text = 'Trained Staff'

            for row_idx, row_data in enumerate(region_df.itertuples(), start=1):
                region_table.rows[row_idx].cells[0].text = str(row_data.Region)
                region_table.rows[row_idx].cells[1].text = str(row_data._2)  # Total Staff
                region_table.rows[row_idx].cells[2].text = str(row_data._3)  # Trained Staff

            # Add chart
            chart_img = create_chart_image(
                region_df[['Region', 'Total Staff', 'Trained Staff']],
                'grouped_bar',
                'Regional Coverage Analysis',
                x='Region', y=['Total Staff', 'Trained Staff']
            )
            doc.add_paragraph()
            doc.add_picture(chart_img, width=DocxInches(6))

            # Add Training by Region bar chart
            if 'Total Trainings' in region_df.columns:
                doc.add_paragraph()
                doc.add_heading('Trainings by Region', 2)
                chart_img2 = create_chart_image(
                    region_df[['Region', 'Total Trainings']].sort_values('Total Trainings', ascending=False),
                    'bar',
                    'Total Trainings by Region',
                    x='Region', y='Total Trainings'
                )
                doc.add_picture(chart_img2, width=DocxInches(6))

    doc.add_page_break()

    # Monthly Trends
    if 'monthly_trends' in report_data and not report_data['monthly_trends'].empty:
        doc.add_heading(f"Monthly Training Trends ({report_data.get('year', date.today().year)})", 1)

        monthly_df = report_data['monthly_trends']

        monthly_table = doc.add_table(rows=len(monthly_df) + 1, cols=3)
        monthly_table.style = 'Light Grid Accent 1'

        monthly_table.rows[0].cells[0].text = 'Month'
        monthly_table.rows[0].cells[1].text = 'Trainings'
        monthly_table.rows[0].cells[2].text = 'Unique Staff'

        for row_idx, row_data in enumerate(monthly_df.itertuples(), start=1):
            monthly_table.rows[row_idx].cells[0].text = str(row_data._4)  # Month Name
            monthly_table.rows[row_idx].cells[1].text = str(row_data.Trainings)
            monthly_table.rows[row_idx].cells[2].text = str(row_data._3)  # Unique Staff

        # Add chart
        chart_img = create_chart_image(
            monthly_df[['Month Name', 'Trainings', 'Unique Staff']],
            'line',
            'Monthly Training Activity',
            x='Month Name', y=['Trainings', 'Unique Staff']
        )
        doc.add_paragraph()
        doc.add_picture(chart_img, width=DocxInches(6))

    doc.add_page_break()

    # Staff Yet to Train
    if 'staff_yet_to_train' in report_data and not report_data['staff_yet_to_train'].empty:
        doc.add_heading('Staff Yet to be Trained', 1)

        yet_df = report_data['staff_yet_to_train']

        doc.add_paragraph(f"Total Staff Needing Training: {len(yet_df)}")

        # Top departments
        if 'department' in yet_df.columns:
            dept_counts = yet_df['department'].value_counts().head(5)
            doc.add_heading('Top Departments Needing Training:', 2)

            for dept, count in dept_counts.items():
                doc.add_paragraph(f"{dept}: {count} staff", style='List Bullet')

    # Save to bytes
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.getvalue()


def create_pdf_report(report_data: dict) -> bytes:
    """
    Creates a PDF report from comprehensive report data.

    Args:
        report_data: Dictionary containing all report sections

    Returns:
        bytes: PDF file as bytes
    """
    bio = io.BytesIO()
    doc = SimpleDocTemplate(bio, pagesize=letter)

    # Container for elements
    elements = []

    # Styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#003366'),
        spaceAfter=30,
        alignment=1  # Center
    )
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#003366'),
        spaceAfter=12
    )

    # Title
    elements.append(Paragraph("Food Concepts Training Report", title_style))
    elements.append(Paragraph(f"Comprehensive Training Analytics<br/>Generated: {date.today().strftime('%B %d, %Y')}", styles['Normal']))
    elements.append(Spacer(1, 0.3*inch))

    # Executive Summary
    elements.append(Paragraph("Executive Summary", heading_style))

    summary_data = report_data.get('summary', {})
    summary_list = [
        ['Metric', 'Value'],
        ['Report Year', str(report_data.get('year', date.today().year))],
        ['Total Trainings', str(summary_data.get('total_trainings', 0))],
        ['Total Staff (All)', str(summary_data.get('total_staff', 0))],
        ['Active Staff', str(summary_data.get('active_staff', 0))],
        ['Staff Trained (Active)', str(summary_data.get('staff_trained', 0))],
        ['Stores Trained', str(summary_data.get('stores_trained', 0))],
        ['Pending Trainings', str(summary_data.get('pending_trainings', 0))],
    ]

    summary_table = Table(summary_list, colWidths=[3*inch, 2*inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))
    elements.append(summary_table)
    elements.append(PageBreak())

    # Training Frequency
    if 'training_frequency' in report_data and not report_data['training_frequency'].empty:
        elements.append(Paragraph("Training Frequency", heading_style))

        freq_df = report_data['training_frequency']
        freq_list = [['Training Title', 'Times Held', 'First Date', 'Last Date']]

        for row in freq_df.itertuples():
            freq_list.append([
                str(row[1]),
                str(row[2]),
                str(row[4]) if len(row) > 4 else "",
                str(row[5]) if len(row) > 5 else ""
            ])

        freq_table = Table(freq_list, colWidths=[2.5*inch, 1*inch, 1.5*inch, 1.5*inch])
        freq_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
        ]))
        elements.append(freq_table)
        elements.append(PageBreak())

    # Stores Coverage
    if 'stores_coverage' in report_data:
        elements.append(Paragraph("Stores Training Coverage", heading_style))

        stores_data = report_data['stores_coverage']
        stores_list = [
            ['Metric', 'Value'],
            ['Total Stores', str(stores_data.get('total', 0))],
            ['Stores Trained', f"{stores_data.get('trained', 0)} ({stores_data.get('trained_pct', 0):.1f}%)"],
            ['Stores Pending', f"{stores_data.get('pending', 0)} ({stores_data.get('pending_pct', 0):.1f}%)"],
        ]

        stores_table = Table(stores_list, colWidths=[3*inch, 2*inch])
        stores_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(stores_table)
        elements.append(Spacer(1, 0.3*inch))

    # Staff Participation
    if 'staff_participation' in report_data:
        elements.append(Paragraph("Staff Participation Analysis", heading_style))

        staff_data = report_data['staff_participation']
        staff_list = [
            ['Metric', 'Value'],
            ['Eligible Staff', str(staff_data.get('eligible', 0))],
            ['Participated', f"{staff_data.get('participated', 0)} ({staff_data.get('participation_rate', 0):.1f}%)"],
            ['Not Participated', str(staff_data.get('not_participated', 0))],
            ['Participation Rate', f"{staff_data.get('participation_rate', 0):.1f}%"],
        ]

        staff_table = Table(staff_list, colWidths=[3*inch, 2*inch])
        staff_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(staff_table)
        elements.append(PageBreak())

    # Food Safety Report
    if 'food_safety' in report_data:
        elements.append(Paragraph("Food Safety Compliance Report", heading_style))

        food_data = report_data['food_safety']
        food_list = [
            ['Metric', 'Value'],
            ['Required to be Tested', str(food_data.get('required', 0))],
            ['Already Tested', f"{food_data.get('tested', 0)} ({food_data.get('tested_pct', 0):.1f}%)"],
            ['Never Tested', str(food_data.get('never_tested', 0))],
            ['Due for Testing', str(food_data.get('due_count', 0))],
        ]

        food_table = Table(food_list, colWidths=[3*inch, 2*inch])
        food_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(food_table)
        elements.append(PageBreak())

    # Compliance Dashboard
    if 'compliance' in report_data and report_data['compliance']:
        elements.append(Paragraph("Compliance Dashboard - All Required Trainings", heading_style))

        compliance_df = pd.DataFrame(report_data['compliance'])
        compliance_list = [['Training', 'Trained Staff', 'Compliance Rate (%)', 'Overdue']]

        for row in compliance_df.itertuples():
            compliance_list.append([
                str(row.Training),
                str(row._2),  # Trained Staff
                f"{row._3:.1f}",  # Compliance Rate
                str(row.Overdue)
            ])

        compliance_table = Table(compliance_list, colWidths=[2.5*inch, 1.2*inch, 1.5*inch, 1*inch])
        compliance_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(compliance_table)

        # Add chart
        chart_img = create_chart_image(
            compliance_df[['Training', 'Compliance Rate']],
            'horizontal_bar',
            'Compliance Rate by Training Type',
            x='Training', y='Compliance Rate'
        )
        elements.append(Spacer(1, 0.2*inch))
        elements.append(Image(chart_img, width=6*inch, height=3.5*inch))
        elements.append(PageBreak())

    # Venue Analysis
    if 'venue_analysis' in report_data and not report_data['venue_analysis'].empty:
        elements.append(Paragraph("Training by Venue Analysis", heading_style))

        venue_df = report_data['venue_analysis']
        venue_list = [['Venue', 'Sessions', 'Unique Staff']]

        for row in venue_df.itertuples():
            venue_list.append([
                str(row.Venue),
                str(row.Sessions),
                str(row._3)  # Unique Staff
            ])

        venue_table = Table(venue_list, colWidths=[3*inch, 1.5*inch, 1.5*inch])
        venue_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(venue_table)

        # Add chart
        chart_img = create_chart_image(
            venue_df[['Venue', 'Sessions']].head(10),
            'bar',
            'Top 10 Training Venues',
            x='Venue', y='Sessions'
        )
        elements.append(Spacer(1, 0.2*inch))
        elements.append(Image(chart_img, width=6*inch, height=3.5*inch))
        elements.append(PageBreak())

    # Department Coverage
    if 'department_coverage' in report_data and not report_data['department_coverage'].empty:
        dept_df = report_data['department_coverage']
        # Filter out rows with zero or negative values
        dept_df = dept_df[(dept_df['Total Staff'] > 0) & (dept_df['Total Trainings'] > 0)]

        if not dept_df.empty:
            elements.append(Paragraph("Department Training Coverage", heading_style))

            dept_list = [['Department', 'Total Staff', 'Trained Staff', 'Coverage Rate (%)']]

            for row in dept_df.itertuples():
                coverage_rate = getattr(row, '_5', 0) if hasattr(row, '_5') else 0
                dept_list.append([
                    str(row.Department),
                    str(row._2),  # Total Staff
                    str(row._3),  # Trained Staff
                    f"{coverage_rate:.1f}" if coverage_rate > 0 else "N/A"
                ])

            dept_table = Table(dept_list, colWidths=[2*inch, 1.5*inch, 1.5*inch, 1.5*inch])
            dept_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            elements.append(dept_table)

            # Add chart
            chart_img = create_chart_image(
                dept_df[['Department', 'Total Staff', 'Trained Staff']],
                'grouped_bar',
                'Department Coverage Analysis',
                x='Department', y=['Total Staff', 'Trained Staff']
            )
            elements.append(Spacer(1, 0.2*inch))
            elements.append(Image(chart_img, width=6*inch, height=3.5*inch))

            # Add Training by Department bar chart
            if 'Total Trainings' in dept_df.columns:
                elements.append(Spacer(1, 0.3*inch))
                elements.append(Paragraph("Trainings by Department", heading_style))
                chart_img2 = create_chart_image(
                    dept_df[['Department', 'Total Trainings']].sort_values('Total Trainings', ascending=False),
                    'bar',
                    'Total Trainings by Department',
                    x='Department', y='Total Trainings'
                )
                elements.append(Spacer(1, 0.2*inch))
                elements.append(Image(chart_img2, width=6*inch, height=3.5*inch))

            elements.append(PageBreak())

    # Regional Coverage
    if 'regional_coverage' in report_data and not report_data['regional_coverage'].empty:
        region_df = report_data['regional_coverage']
        # Filter out rows with zero or negative values
        region_df = region_df[(region_df['Total Staff'] > 0) & (region_df['Total Trainings'] > 0)]

        if not region_df.empty:
            elements.append(Paragraph("Regional Training Coverage", heading_style))

            region_list = [['Region', 'Total Staff', 'Trained Staff']]

            for row in region_df.itertuples():
                region_list.append([
                    str(row.Region),
                    str(row._2),  # Total Staff
                    str(row._3)  # Trained Staff
                ])

            region_table = Table(region_list, colWidths=[2.5*inch, 1.5*inch, 1.5*inch])
            region_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            elements.append(region_table)

            # Add chart
            chart_img = create_chart_image(
                region_df[['Region', 'Total Staff', 'Trained Staff']],
                'grouped_bar',
                'Regional Coverage Analysis',
                x='Region', y=['Total Staff', 'Trained Staff']
            )
            elements.append(Spacer(1, 0.2*inch))
            elements.append(Image(chart_img, width=6*inch, height=3.5*inch))

            # Add Training by Region bar chart
            if 'Total Trainings' in region_df.columns:
                elements.append(Spacer(1, 0.3*inch))
                elements.append(Paragraph("Trainings by Region", heading_style))
                chart_img2 = create_chart_image(
                    region_df[['Region', 'Total Trainings']].sort_values('Total Trainings', ascending=False),
                    'bar',
                    'Total Trainings by Region',
                    x='Region', y='Total Trainings'
                )
                elements.append(Spacer(1, 0.2*inch))
                elements.append(Image(chart_img2, width=6*inch, height=3.5*inch))

            elements.append(PageBreak())

    # Monthly Trends
    if 'monthly_trends' in report_data and not report_data['monthly_trends'].empty:
        elements.append(Paragraph(f"Monthly Training Trends ({report_data.get('year', date.today().year)})", heading_style))

        monthly_df = report_data['monthly_trends']
        monthly_list = [['Month', 'Trainings', 'Unique Staff']]

        for row in monthly_df.itertuples():
            monthly_list.append([
                str(row._4),  # Month Name
                str(row.Trainings),
                str(row._3)  # Unique Staff
            ])

        monthly_table = Table(monthly_list, colWidths=[2*inch, 2*inch, 2*inch])
        monthly_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(monthly_table)

        # Add chart
        chart_img = create_chart_image(
            monthly_df[['Month Name', 'Trainings', 'Unique Staff']],
            'line',
            'Monthly Training Activity',
            x='Month Name', y=['Trainings', 'Unique Staff']
        )
        elements.append(Spacer(1, 0.2*inch))
        elements.append(Image(chart_img, width=6*inch, height=3.5*inch))
        elements.append(PageBreak())

    # Staff Yet to Train
    if 'staff_yet_to_train' in report_data and not report_data['staff_yet_to_train'].empty:
        elements.append(Paragraph("Staff Yet to be Trained", heading_style))

        yet_df = report_data['staff_yet_to_train']

        elements.append(Paragraph(f"Total Staff Needing Training: {len(yet_df)}", styles['Normal']))
        elements.append(Spacer(1, 0.2*inch))

        # Top departments
        if 'department' in yet_df.columns:
            dept_counts = yet_df['department'].value_counts().head(5)

            elements.append(Paragraph("Top Departments Needing Training:", styles['Heading3']))

            dept_list = [['Department', 'Staff Count']]
            for dept, count in dept_counts.items():
                dept_list.append([str(dept), str(count)])

            dept_table = Table(dept_list, colWidths=[3*inch, 2*inch])
            dept_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#003366')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            elements.append(dept_table)

    # Build PDF
    doc.build(elements)
    bio.seek(0)
    return bio.getvalue()
