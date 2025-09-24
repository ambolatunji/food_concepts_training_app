import io, os, json, subprocess
from pathlib import Path
from datetime import date, datetime

import pandas as pd
import plotly.express as px
import streamlit as st

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from core.db import (
    get_conn, list_trainings, count_employees, count_employees_matching,
    count_trained_employees, employees_due_flags
)

# -------------- Page setup
st.set_page_config(page_title="Reporting", page_icon="📑", layout="wide")
st.title("📑 Reporting")

conn = get_conn()

# -------------- Helpers
def _chart_to_png(fig):
    # Requires kaleido
    try:
        import plotly.io as pio
        return pio.to_image(fig, format="png", width=1280, height=720, scale=2)
    except Exception:
        return None

def _insert_image_slide(prs, title, png_bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    left = Inches(0.5); top = Inches(0.5)
    if title:
        tx = slide.shapes.add_textbox(left, top, Inches(12), Inches(0.6))
        tf = tx.text_frame; tf.text = title
        tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        top = Inches(1.2)
    if png_bytes:
        slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=Inches(12))

def _ppt_infographic_kpis(prs, metrics):
    """
    Adds a slide with 6 KPI 'cards' (rounded rectangles).
    metrics: list of dicts with keys: title, value, sub (optional), color (RGB tuple)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.6))
    t = title.text_frame; t.text = "Key Performance Indicators"
    t.paragraphs[0].font.size = Pt(26)

    # layout grid 3 x 2
    card_w, card_h = Inches(3.9), Inches(1.5)
    x0, y0 = Inches(0.5), Inches(1.2)
    gap_x, gap_y = Inches(0.3), Inches(0.3)

    def add_card(i, m):
        row, col = divmod(i, 3)
        left = x0 + col*(card_w + gap_x)
        top  = y0 + row*(card_h + gap_y)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        r, g, b = m.get("color", (38, 140, 222))
        fill = shape.fill
        fill.solid(); fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.color.rgb = RGBColor(255, 255, 255)

        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = m["title"]
        p1.font.size = Pt(12); p1.font.bold = True; p1.font.color.rgb = RGBColor(255,255,255)

        p2 = tf.add_paragraph()
        p2.text = str(m["value"])
        p2.font.size = Pt(28); p2.font.bold = True; p2.font.color.rgb = RGBColor(255,255,255)

        sub = m.get("sub")
        if sub:
            p3 = tf.add_paragraph()
            p3.text = sub
            p3.font.size = Pt(10); p3.font.color.rgb = RGBColor(255,255,255)

    for i, m in enumerate(metrics[:6]):
        add_card(i, m)

def _ppt_table(slide, top_left, data, col_widths=None, title=None):
    """
    Add a simple table to the given slide. data: list of rows (first row = headers)
    top_left: (left_inches, top_inches)
    col_widths: list of inches for each column
    """
    rows = len(data)
    cols = len(data[0]) if rows else 0
    if rows == 0 or cols == 0:
        return

    left = Inches(top_left[0]); top = Inches(top_left[1])
    width = Inches(sum(col_widths) if col_widths else 12)
    height = Inches(0.8 + 0.3*rows)

    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            tbl.cell(r, c).text = str(data[r][c])

    # Bold header row
    for c in range(cols):
        tbl.cell(0, c).text_frame.paragraphs[0].font.bold = True

    if title:
        tx = slide.shapes.add_textbox(left, Inches(top_left[1]-0.6), Inches(12), Inches(0.5))
        tx.text_frame.text = title
        tx.text_frame.paragraphs[0].font.size = Pt(22)

def _save_ppt(prs, filename)->str:
    p = Path("data")/filename
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(p.as_posix())
    return p.as_posix()

def _to_pdf_via_libreoffice(path:str)->str:
    outdir = str(Path(path).parent)
    try:
        subprocess.run(["soffice","--headless","--convert-to","pdf","--outdir", outdir, path],
                       check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        pdf_path = path.replace(".pptx",".pdf") if path.endswith(".pptx") else path.replace(".docx",".pdf")
        return pdf_path if Path(pdf_path).exists() else ""
    except Exception:
        return ""
# --- Image helpers: Plotly -> PNG (kaleido) OR fallback to Matplotlib
def _chart_to_png(fig):
    """Try Plotly+kaleido; return PNG bytes or None."""
    try:
        import plotly.io as pio
        return pio.to_image(fig, format="png", width=1280, height=720, scale=2)
    except Exception:
        return None

def _to_png_matplotlib(kind, df, *, x=None, y=None, hue=None, title="", rotate=False, stacked=False, cmap="Blues"):
    """Minimal Matplotlib fallback. Returns PNG bytes."""
    import io as _io
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=(12, 6))

    if kind == "bar":
        if hue and stacked:
            pivot = df.pivot(index=x, columns=hue, values=y).fillna(0)
            pivot.plot(kind="bar", stacked=True, ax=ax)
        else:
            ax.bar(df[x], df[y])
    elif kind == "line":
        ax.plot(df[x], df[y], marker="o")
    elif kind == "lines_multi":
        for col in y:  # y is a list
            ax.plot(df[x], df[col], marker="o", label=col)
        ax.legend()
    elif kind == "pie":
        ax.pie(df[y], labels=df[x], autopct="%1.0f%%", startangle=140)
        ax.axis("equal")
    elif kind == "heatmap":
        import numpy as np
        im = ax.imshow(df.values, aspect="auto", cmap=cmap)
        ax.set_xticks(range(df.shape[1])); ax.set_xticklabels(df.columns, rotation=45, ha="right")
        ax.set_yticks(range(df.shape[0])); ax.set_yticklabels(df.index)
        fig.colorbar(im, ax=ax)

    ax.set_title(title)
    if rotate:
        plt.setp(ax.get_xticklabels(), rotation=45, ha="right")
    buf = _io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="png", dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


# -------- OpenRouter client (from your previous fix)
import requests
OPENROUTER_BASE = "https://openrouter.ai/api/v1"
def call_openrouter(model: str, messages: list, api_key: str, temperature: float = 0.3, referer: str = "http://localhost", title: str = "Food Concepts Training Report") -> str:
    if not api_key or not api_key.strip():
        raise RuntimeError("Missing OpenRouter API key.")
    headers = {
        "Authorization": f"Bearer {api_key.strip()}",
        "HTTP-Referer": referer,
        "X-Title": title,
        "Content-Type": "application/json",
    }
    payload = {"model": model, "messages": messages, "temperature": temperature}
    url_chat = f"{OPENROUTER_BASE}/chat/completions"
    r = requests.post(url_chat, headers=headers, json=payload, timeout=120)
    if r.status_code == 200:
        j = r.json()
        return j["choices"][0]["message"]["content"]
    # fallback legacy /completions on 404
    if r.status_code == 404:
        prompt = "\n\n".join([m.get("content","") for m in messages])
        r2 = requests.post(f"{OPENROUTER_BASE}/completions", headers=headers, json={"model": payload["model"], "prompt": prompt, "temperature": temperature}, timeout=120)
        r2.raise_for_status()
        j2 = r2.json()
        return j2["choices"][0]["text"]
    r.raise_for_status()
    return ""

# -------------- Data
rows = list_trainings(conn, {})
df = pd.DataFrame(rows, columns=[
    "Training ID","Employee Name","Department","Store","Position","Region",
    "Training Date","Next Due Date","Evidence Path","Training Title","Training Venue"
])

total_emps = count_employees(conn)
trained_emps = count_trained_employees(conn, {})
flags = employees_due_flags(conn, {})
due_30 = sum(1 for v in flags.values() if v["any_due_30"])
overdue = sum(1 for v in flags.values() if v["any_overdue"])
coverage = (trained_emps/total_emps*100.0) if total_emps else 0.0

# enrich df
if not df.empty:
    df["Training Date"] = pd.to_datetime(df["Training Date"], errors="coerce")
    df["Next Due Date"] = pd.to_datetime(df["Next Due Date"], errors="coerce")

# -------------- Derived analytics & charts
def chart_trainings_by_department(df):
    if df.empty: return None
    g = df.groupby("Department")["Training ID"].count().reset_index(name="Trainings").sort_values("Trainings", ascending=False)
    fig = px.bar(g, x="Department", y="Trainings", title="Trainings by Department")
    fig.update_layout(margin=dict(l=30,r=10,b=120,t=60), xaxis_tickangle=-45)
    png = _chart_to_png(fig)
    return png or _to_png_matplotlib("bar", g, x="Department", y="Trainings", title="Trainings by Department", rotate=True)

def chart_trainings_by_title(df):
    if df.empty: return None
    g = df.groupby("Training Title")["Training ID"].count().reset_index(name="Trainings").sort_values("Trainings", ascending=False)
    fig = px.bar(g, x="Training Title", y="Trainings", title="Trainings by Title")
    fig.update_layout(margin=dict(l=30,r=10,b=160,t=60), xaxis_tickangle=-45)
    png = _chart_to_png(fig)
    return png or _to_png_matplotlib("bar", g, x="Training Title", y="Trainings", title="Trainings by Title", rotate=True)

def chart_trainings_by_month(df):
    if df.empty: return None
    d = df.dropna(subset=["Training Date"]).copy()
    if d.empty: return None
    d["Month"] = d["Training Date"].dt.to_period("M").astype(str)
    g = d.groupby("Month")["Training ID"].count().reset_index(name="Trainings")
    fig = px.line(g, x="Month", y="Trainings", title="Training Trend by Month", markers=True)
    fig.update_layout(margin=dict(l=30,r=10,b=80,t=60))
    png = _chart_to_png(fig)
    return png or _to_png_matplotlib("line", g, x="Month", y="Trainings", title="Training Trend by Month")

def coverage_by_department(conn, df):
    emps = pd.read_sql_query("SELECT id, department FROM employees WHERE deleted_at IS NULL", conn)
    if emps.empty:
        return pd.DataFrame(columns=["Department","Headcount","Trained","Coverage"]), None

    head = emps.groupby("department")["id"].nunique().reset_index(name="Headcount")
    if df.empty:
        cov = head.assign(Trained=0, Coverage=0.0)
    else:
        trained = df.groupby("Department")["Employee Name"].nunique().reset_index(name="Trained")
        cov = head.merge(trained, left_on="department", right_on="Department", how="left").fillna({"Trained":0})
        cov["Coverage"] = (cov["Trained"]/cov["Headcount"]*100.0).round(1)
    cov.rename(columns={"department":"Department"}, inplace=True)

    if cov.empty: return cov, None
    fig = px.bar(cov.sort_values("Coverage", ascending=False), x="Department", y="Coverage", title="Coverage by Department (%)")
    fig.update_layout(margin=dict(l=30,r=10,b=160,t=60), xaxis_tickangle=-45, yaxis_range=[0,100])
    png = _chart_to_png(fig)
    return cov, (png or _to_png_matplotlib("bar", cov.sort_values("Coverage", ascending=False),
                                           x="Department", y="Coverage", title="Coverage by Department (%)", rotate=True))

def coverage_treemap_by_store(conn, df):
    # Treemap (Plotly) -> bar fallback
    emps = pd.read_sql_query("SELECT id, region, store FROM employees WHERE deleted_at IS NULL", conn)
    if emps.empty: return None
    head = emps.groupby(["region","store"])["id"].nunique().reset_index(name="Headcount")
    if df.empty:
        cov = head.assign(Trained=0, Coverage=0.0)
    else:
        trained = df.groupby("Store")["Employee Name"].nunique().reset_index(name="Trained")
        cov = head.merge(trained, left_on="store", right_on="Store", how="left").fillna({"Trained":0})
        cov["Coverage"] = (cov["Trained"]/cov["Headcount"]*100.0).round(1)
    cov.rename(columns={"region":"Region","store":"Store"}, inplace=True)
    if cov.empty: return None

    fig = px.treemap(cov, path=["Region","Store"], values="Headcount",
                     color="Coverage", color_continuous_scale="RdYlGn",
                     title="Store Coverage Treemap (size=headcount, color=coverage%)")
    fig.update_layout(margin=dict(l=10,r=10,b=10,t=60))
    png = _chart_to_png(fig)
    if png: return png

    # Fallback: horizontal bar by coverage (top 30)
    cov2 = cov.sort_values("Coverage", ascending=False).head(30)
    return _to_png_matplotlib("bar", cov2, x="Store", y="Coverage",
                              title="Store Coverage (Top 30 by %)", rotate=True)

def due_overdue_by_department(conn, df):
    emps = pd.read_sql_query("SELECT id, department FROM employees WHERE deleted_at IS NULL", conn)
    if emps.empty: return None, None
    fl = employees_due_flags(conn, {})
    fdf = pd.DataFrame([{"id":k, **v} for k,v in fl.items()]) if fl else pd.DataFrame(columns=["id","any_due_30","any_overdue"])
    base = emps.merge(fdf, on="id", how="left").fillna({"any_due_30":0,"any_overdue":0})
    g = base.groupby("department")[["any_due_30","any_overdue"]].sum().reset_index()
    g.rename(columns={"department":"Department","any_due_30":"Due≤30d","any_overdue":"Overdue"}, inplace=True)

    fig = px.bar(g.sort_values("Due≤30d", ascending=False), x="Department", y=["Due≤30d","Overdue"],
                 barmode="stack", title="Due vs Overdue by Department")
    fig.update_layout(margin=dict(l=30,r=10,b=160,t=60), xaxis_tickangle=-45)
    png = _chart_to_png(fig)
    if png: return g, png
    # Matplotlib stacked
    g_m = g.melt("Department", ["Due≤30d","Overdue"], var_name="Status", value_name="Count")
    return g, _to_png_matplotlib("bar", g_m, x="Department", y="Count", hue="Status",
                                 title="Due vs Overdue by Department", rotate=True, stacked=True)

def top_stores_by_trainings(df, topn=20):
    if df.empty: return None
    g = df.groupby("Store")["Training ID"].count().reset_index(name="Trainings").sort_values("Trainings", ascending=False).head(topn)
    fig = px.bar(g, x="Store", y="Trainings", title=f"Top {topn} Stores by Trainings")
    fig.update_layout(margin=dict(l=30,r=10,b=200,t=60), xaxis_tickangle=-65)
    png = _chart_to_png(fig)
    return png or _to_png_matplotlib("bar", g, x="Store", y="Trainings", title=f"Top {topn} Stores by Trainings", rotate=True)

def joins_vs_leaves_by_month(conn):
    emp = pd.read_sql_query("SELECT start_date, end_date FROM employees WHERE deleted_at IS NULL", conn)
    if emp.empty: return None, None
    emp["start_date"] = pd.to_datetime(emp["start_date"], errors="coerce")
    emp["end_date"]   = pd.to_datetime(emp["end_date"], errors="coerce")
    j = emp.dropna(subset=["start_date"]).assign(Month=lambda d: d["start_date"].dt.to_period("M").astype(str)).groupby("Month").size().reset_index(name="Joins")
    l = emp.dropna(subset=["end_date"]).assign(Month=lambda d: d["end_date"].dt.to_period("M").astype(str)).groupby("Month").size().reset_index(name="Leaves")
    g = pd.merge(j, l, on="Month", how="outer").fillna(0).sort_values("Month")

    fig = px.bar(g, x="Month", y=["Joins","Leaves"], barmode="group", title="Joins vs Leaves by Month")
    fig.update_layout(margin=dict(l=30,r=10,b=80,t=60))
    png = _chart_to_png(fig)
    if png: return g, png
    # Matplotlib line fallback (two series)
    return g, _to_png_matplotlib("lines_multi", g, x="Month", y=["Joins","Leaves"], title="Joins vs Leaves by Month")


def to_docx(report_text:str, charts:dict, filename:str="Training_Report_Detail.docx")->str:
    from docx import Document
    from docx.shared import Inches, Pt
    doc = Document()
    doc.add_heading(f"Food Concepts – Training Report ({date.today().isoformat()})", 0)

    # Executive summary
    doc.add_heading("Executive Summary", level=1)
    for para in report_text.split("\n"):
        if para.strip():
            doc.add_paragraph(para.strip())

    # KPIs
    doc.add_heading("Key Metrics", level=1)
    doc.add_paragraph(f"Total Employees: {total_emps}")
    doc.add_paragraph(f"Employees Trained (≥1): {trained_emps} ({coverage:.1f}%)")
    doc.add_paragraph(f"Due ≤30 days: {due_30}")
    doc.add_paragraph(f"Overdue: {overdue}")

    # Charts
    for title, png in charts.items():
        if not png: 
            continue
        doc.add_heading(title, level=2)
        tmp = Path("data")/f"{title.replace(' ','_')}.png"
        tmp.write_bytes(png)
        doc.add_picture(tmp.as_posix(), width=Inches(6.5))

    out = Path("data")/filename
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(out.as_posix())
    return out.as_posix()

def region_pie_by_trainings(df):
    if df.empty: return None
    g = df.groupby("Region")["Training ID"].count().reset_index(name="Trainings")
    if g.empty: return None
    fig = px.pie(g, names="Region", values="Trainings", title="Trainings Share by Region", hole=0.35)
    fig.update_layout(margin=dict(l=10,r=10,b=10,t=60))
    png = _chart_to_png(fig)
    return png or _to_png_matplotlib("pie", g, x="Region", y="Trainings", title="Trainings Share by Region")

def heatmap_trainings_by_dept_month(df):
    if df.empty or "Training Date" not in df.columns: return None
    d = df.dropna(subset=["Training Date"]).copy()
    if d.empty: return None
    d["Month"] = d["Training Date"].dt.to_period("M").astype(str)
    pvt = d.pivot_table(index="Department", columns="Month", values="Training ID", aggfunc="count", fill_value=0)
    if pvt.empty: return None
    fig = px.imshow(pvt, aspect="auto", title="Heatmap: Trainings by Department × Month", color_continuous_scale="Blues")
    fig.update_layout(margin=dict(l=40,r=20,b=120,t=60), xaxis_tickangle=-45)
    png = _chart_to_png(fig)
    if png: return png
    return _to_png_matplotlib("heatmap", pvt, title="Heatmap: Trainings by Department × Month", cmap="Blues")


def stacked_titles_by_department(df, top_titles=6):
    if df.empty: return None
    top = df["Training Title"].value_counts().head(top_titles).index.tolist()
    d = df[df["Training Title"].isin(top)]
    if d.empty: return None
    g = d.groupby(["Department","Training Title"])["Training ID"].count().reset_index(name="Trainings")
    fig = px.bar(g, x="Department", y="Trainings", color="Training Title", barmode="stack",
                 title=f"Stacked Trainings by Department (Top {top_titles} Titles)")
    fig.update_layout(margin=dict(l=30,r=10,b=160,t=60), xaxis_tickangle=-45)
    png = _chart_to_png(fig)
    if png: return png
    # Matplotlib stacked
    return _to_png_matplotlib("bar", g, x="Department", y="Trainings", hue="Training Title",
                              title=f"Stacked Trainings by Department (Top {top_titles} Titles)", rotate=True, stacked=True)


def coverage_treemap_by_store(conn, df):
    # Headcount & trained per store, color by coverage
    emps = pd.read_sql_query("SELECT id, region, store FROM employees WHERE deleted_at IS NULL", conn)
    if emps.empty: return None
    head = emps.groupby(["region","store"])["id"].nunique().reset_index(name="Headcount")
    if df.empty:
        cov = head.assign(Trained=0, Coverage=0.0)
    else:
        trained = df.groupby("Store")["Employee Name"].nunique().reset_index(name="Trained")
        cov = head.merge(trained, left_on="store", right_on="Store", how="left").fillna({"Trained":0})
        cov["Coverage"] = (cov["Trained"]/cov["Headcount"]*100.0).round(1)
    cov.rename(columns={"region":"Region","store":"Store"}, inplace=True)
    if cov.empty: return None
    fig = px.treemap(cov, path=["Region","Store"], values="Headcount",
                     color="Coverage", color_continuous_scale="RdYlGn",
                     title="Store Coverage Treemap (size=headcount, color=coverage%)")
    fig.update_layout(margin=dict(l=10,r=10,b=10,t=60))
    return _chart_to_png(fig)

def turnover_by_department(conn):
    emp = pd.read_sql_query("SELECT department, end_date FROM employees WHERE deleted_at IS NULL", conn)
    if emp.empty: return None
    emp["end_date"] = pd.to_datetime(emp["end_date"], errors="coerce")
    # simple period = last 12 months
    cutoff = pd.Timestamp.today().normalize() - pd.offsets.DateOffset(years=1)
    leavers = emp[emp["end_date"] >= cutoff]
    if leavers.empty: return None
    g = leavers.groupby("department").size().reset_index(name="Leavers (12m)")
    # approximate headcount per dept now
    head = pd.read_sql_query("SELECT department, COUNT(*) as Headcount FROM employees WHERE deleted_at IS NULL GROUP BY department", conn)
    cov = g.merge(head, on="department", how="left")
    cov["Turnover % (12m)"] = (cov["Leavers (12m)"] / cov["Headcount"] * 100.0).round(1)
    cov.rename(columns={"department":"Department"}, inplace=True)
    fig = px.bar(cov.sort_values("Turnover % (12m)", ascending=False), x="Department", y="Turnover % (12m)",
                 title="Turnover by Department (last 12 months)")
    fig.update_layout(margin=dict(l=30,r=10,b=160,t=60), xaxis_tickangle=-45, yaxis_range=[0, cov["Turnover % (12m)"].max()*1.2])
    return _chart_to_png(fig)

def get_secret(key:str, default:str="")->str:
    # from secrets.toml if available (Streamlit Cloud)
    try:
        return st.secrets[key]
    except Exception:
        return os.getenv(key.upper(), default)
# -------------- Tabs
tab1, tab2 = st.tabs(["🤖 AI-generated report", "🧭 Deterministic report"])

# ===== Common charts to reuse in both flows =====
png_by_dept   = chart_trainings_by_department(df)
png_by_title  = chart_trainings_by_title(df)
png_by_month  = chart_trainings_by_month(df)
cov_df, png_cov_dept = coverage_by_department(conn, df)
due_df, png_due_dept = due_overdue_by_department(conn, df)
png_top_stores = top_stores_by_trainings(df)
jl_df, png_jl  = joins_vs_leaves_by_month(conn)
png_region_pie  = region_pie_by_trainings(df)
png_heat_dept   = heatmap_trainings_by_dept_month(df)
png_stack_dept  = stacked_titles_by_department(df, top_titles=6)
png_cov_tree    = coverage_treemap_by_store(conn, df)
png_turnover_dept = turnover_by_department(conn)


# ---------- AI Report ----------
with tab1:
    st.subheader("🤖 Generate with OpenRouter")
    model = st.selectbox("Model", [
        "meta-llama/llama-3.1-405b-instruct",
        "nvidia/nemotron-nano-9b-v2:free",
        "mistralai/mixtral-8x22b-instruct",
        "openrouter/auto"
    ], index=0)

    #get api key from secrets.toml if available
    api_key = get_secret("openrouter_api_key")
    if not api_key:
        api_key = st.text_input("OpenRouter API Key", type="password")
    if not api_key:
        st.error("Please provide an OpenRouter API key.")
    #api_key = st.text_input("OpenRouter API Key", type="password")
    tpl = st.text_input("PPTX template (optional)", value="data/report_template.pptx")

    if st.button("Build AI Report", type="primary", use_container_width=True):
        # Build summary via OpenRouter (longer prompt for richer slide 2 content)
        sample = df.head(500).to_dict(orient="records")
        try:
            summary = call_openrouter(
                model=model,
                api_key=api_key,
                referer="http://localhost",
                title="Food Concepts Training Report",
                messages=[{
                    "role": "user",
                    "content": f"""
You are a senior analyst. Create a detailed executive summary for a Training & Safety report.
Please include:
- Highlights, risks, and recommended actions
- Which departments/stores need attention and why
- Observations about overdue vs due-soon and training coverage
- Notes on joiners/leavers impact this period
- 3-5 bullet recommendations with measurable next steps

Key metrics:
- Total Employees: {total_emps}
- Employees trained (>=1): {trained_emps} ({coverage:.1f}%)
- Due ≤30d: {due_30}, Overdue: {overdue}

Data sample (JSON list of rows): {json.dumps(sample)[:70000]}
""".strip()
                }]
            )
        except Exception as e:
            st.error(f"OpenRouter error: {e}")
            st.stop()

        # --- Build PPTX with more slides & infographics
        prs = Presentation(tpl) if Path(tpl).exists() else Presentation()

        # Title
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(1.0))
        tf = tx.text_frame; tf.text = f"Food Concepts – Training Report ({date.today().isoformat()})"
        tf.paragraphs[0].font.size = Pt(28)

        # KPI infographic slide
        kpis = [
            {"title":"Total Employees","value": total_emps, "color": (0, 92, 175)},
            {"title":"Employees Trained (≥1)","value": trained_emps, "sub": f"{coverage:.1f}% coverage", "color": (0, 150, 136)},
            {"title":"Due ≤ 30 days","value": due_30, "color": (255, 152, 0)},
            {"title":"Overdue","value": overdue, "color": (230, 74, 25)},
            {"title":"Total Trainings","value": len(df), "color": (63, 81, 181)},
            {"title":"Data Date","value": date.today().isoformat(), "color": (96, 125, 139)},
        ]
        _ppt_infographic_kpis(prs, kpis)

        # AI Summary slide (richer)
        sum_slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = sum_slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(6))
        box.text_frame.text = summary

        # Charts slides
        _insert_image_slide(prs, "Trainings by Department", png_by_dept)
        _insert_image_slide(prs, "Trainings by Title", png_by_title)
        _insert_image_slide(prs, "Training Trend by Month", png_by_month)
        _insert_image_slide(prs, "Coverage by Department (%)", png_cov_dept)
        _insert_image_slide(prs, "Due vs Overdue by Department", png_due_dept)
        _insert_image_slide(prs, "Top Stores by Trainings", png_top_stores)
        _insert_image_slide(prs, "Joins vs Leaves by Month", png_jl)
        _insert_image_slide(prs, "Trainings Share by Region", png_region_pie)
        _insert_image_slide(prs, "Heatmap: Trainings by Department × Month", png_heat_dept)
        _insert_image_slide(prs, "Stacked Trainings by Department (Top Titles)", png_stack_dept)
        _insert_image_slide(prs, "Store Coverage Treemap", png_cov_tree)
        _insert_image_slide(prs, "Turnover by Department (last 12 months)", png_turnover_dept)


        # Due soon table (top 20)
        due_table = []
        if not df.empty:
            soon = df.dropna(subset=["Next Due Date"]).copy()
            soon = soon[soon["Next Due Date"] <= (pd.Timestamp.today().normalize()+pd.Timedelta(days=30))]
            soon = soon.sort_values("Next Due Date").head(20)
            if not soon.empty:
                due_table = [["Employee","Department","Store","Next Due","Title"]]
                for _, r in soon.iterrows():
                    due_table.append([r["Employee Name"], r["Department"], r["Store"], r["Next Due Date"].date().isoformat(), r["Training Title"]])

        if due_table:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            _ppt_table(slide, (0.5, 1.2), due_table, title="Due ≤ 30 days (Top 20)", col_widths=[2.6,2.2,2.2,1.5,3.0])

        ppt_path = _save_ppt(prs, f"AI_Report_{date.today().isoformat()}.pptx")
        st.success(f"PPTX created: {ppt_path}")
        st.download_button("Download PPTX", data=open(ppt_path,"rb").read(),
                           file_name=Path(ppt_path).name,
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           use_container_width=True)

        pdf_path = _to_pdf_via_libreoffice(ppt_path)
        if pdf_path:
            st.download_button("Download PDF", data=open(pdf_path,"rb").read(),
                               file_name=Path(pdf_path).name, mime="application/pdf",
                               use_container_width=True)
        else:
            st.info("Install LibreOffice for PPTX→PDF export (soffice).")

        # --- Build detailed Word document with charts embedded
        charts = {
            "Trainings by Department": png_by_dept,
            "Trainings by Title": png_by_title,
            "Training Trend by Month": png_by_month,
            "Coverage by Department": png_cov_dept,
            "Due vs Overdue by Department": png_due_dept,
            "Top Stores by Trainings": png_top_stores,
            "Joins vs Leaves by Month": png_jl,
            "Trainings Share by Region": png_region_pie,
            "Heatmap: Trainings by Department × Month": png_heat_dept,
            "Stacked Trainings by Department (Top Titles)": png_stack_dept,
            "Store Coverage Treemap": png_cov_tree,
            "Turnover by Department (last 12 months)": png_turnover_dept,
        }
        doc_path = to_docx(summary, charts, filename=f"Training_Report_Detail_{date.today().isoformat()}.docx")
        st.download_button("Download Word Report (DOCX)", data=open(doc_path,"rb").read(),
                           file_name=Path(doc_path).name, mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                           use_container_width=True)

        pdf_doc = _to_pdf_via_libreoffice(doc_path)
        if pdf_doc:
            st.download_button("Download Word Report as PDF", data=open(pdf_doc,"rb").read(),
                               file_name=Path(pdf_doc).name, mime="application/pdf",
                               use_container_width=True)

# ---------- Deterministic Report ----------
with tab2:
    st.subheader("🧭 Generate deterministic PPT from template (no AI)")
    tpl2 = st.text_input("PPTX template (optional)", value="data/report_template.pptx", key="tpl2")
    if st.button("Build Deterministic Report", use_container_width=True):
        prs = Presentation(tpl2) if Path(tpl2).exists() else Presentation()

        # Title
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(1.0))
        tx.text_frame.text = f"Food Concepts – Training Report (Deterministic) – {date.today().isoformat()}"

        # KPI infographic
        kpis = [
            {"title":"Total Employees","value": total_emps, "color": (0, 92, 175)},
            {"title":"Employees Trained (≥1)","value": trained_emps, "sub": f"{coverage:.1f}% coverage", "color": (0, 150, 136)},
            {"title":"Due ≤ 30 days","value": due_30, "color": (255, 152, 0)},
            {"title":"Overdue","value": overdue, "color": (230, 74, 25)},
            {"title":"Total Trainings","value": len(df), "color": (63, 81, 181)},
            {"title":"Data Date","value": date.today().isoformat(), "color": (96, 125, 139)},
        ]
        _ppt_infographic_kpis(prs, kpis)

        # Charts (same set)
        _insert_image_slide(prs, "Trainings by Department", png_by_dept)
        _insert_image_slide(prs, "Trainings by Title", png_by_title)
        _insert_image_slide(prs, "Training Trend by Month", png_by_month)
        _insert_image_slide(prs, "Coverage by Department (%)", png_cov_dept)
        _insert_image_slide(prs, "Due vs Overdue by Department", png_due_dept)
        _insert_image_slide(prs, "Top Stores by Trainings", png_top_stores)
        _insert_image_slide(prs, "Joins vs Leaves by Month", png_jl)

        path = _save_ppt(prs, f"Deterministic_Report_{date.today().isoformat()}.pptx")
        st.success(f"PPTX created: {path}")
        st.download_button("Download PPTX", data=open(path,"rb").read(),
                           file_name=Path(path).name,
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           use_container_width=True)
        pdf = _to_pdf_via_libreoffice(path)
        if pdf:
            st.download_button("Download PDF", data=open(pdf,"rb").read(),
                               file_name=Path(pdf).name, mime="application/pdf",
                               use_container_width=True)
        else:
            st.info("Install LibreOffice for PPTX→PDF export.")
